from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap
from pptx.oxml import parse_xml

# 创建 16:9 宽屏演示文稿
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 配色方案
DARK_BG = RGBColor(0x0a, 0x0e, 0x27)  # 深蓝黑
ACCENT_CYAN = RGBColor(0x00, 0xe5, 0xff)  # 青色
ACCENT_PURPLE = RGBColor(0x7c, 0x4d, 0xff)  # 紫色
TEXT_MAIN = RGBColor(0xe0, 0xf7, 0xfa)  # 主文字
TEXT_SUB = RGBColor(0x90, 0xa4, 0xae)  # 次要文字
WHITE = RGBColor(0xff, 0xff, 0xff)

def add_background(slide, color=DARK_BG):
    """为幻灯片添加纯色背景"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = color
    background.line.fill.background()
    # 将背景移到最底层
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_title(slide, title, subtitle=None, is_cover=False):
    """添加标题"""
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.0))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(36 if is_cover else 28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_CYAN
    p.font.name = "Microsoft YaHei"
    
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(16)
        p2.font.color.rgb = TEXT_SUB
        p2.font.name = "Microsoft YaHei"
        p2.space_before = Pt(8)

def add_bullets(slide, bullets, left=Inches(0.9), top=Inches(1.8), width=Inches(11.5), font_size=18, color=TEXT_MAIN):
    """添加 bullet points"""
    box = slide.shapes.add_textbox(left, top, width, Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(14)
        p.level = 0
    return box

def add_decorative_line(slide, top):
    """添加装饰线"""
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), top, Inches(2.0), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_CYAN
    line.line.fill.background()

def add_footer(slide, text="CCDC AI Service · Vibe Coding 实战"):
    """添加页脚"""
    footer = slide.shapes.add_textbox(Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.3))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_SUB
    p.font.name = "Microsoft YaHei"

# ========== 第1页：封面 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
# 装饰图形
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9.5), Inches(1.5), Inches(3.5), Inches(3.5))
circle.fill.solid()
circle.fill.fore_color.rgb = ACCENT_PURPLE
circle.fill.transparency = 0.85
circle.line.color.rgb = ACCENT_PURPLE
circle.line.width = Pt(1)

rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.0), Inches(2.2), Inches(8.5), Inches(3.5))
rect.fill.background()
rect.line.fill.background()
tf = rect.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Vibe Coding 实战"
p.font.size = Pt(54)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Microsoft YaHei"
p2 = tf.add_paragraph()
p2.text = "零代码手写搭建 CCDC AI 服务平台"
p2.font.size = Pt(28)
p2.font.color.rgb = ACCENT_CYAN
p2.font.name = "Microsoft YaHei"
p2.space_before = Pt(12)
p3 = tf.add_paragraph()
p3.text = "从想法到上线的 AI 应用开发新范式"
p3.font.size = Pt(18)
p3.font.color.rgb = TEXT_SUB
p3.font.name = "Microsoft YaHei"
p3.space_before = Pt(20)

add_footer(slide)

# ========== 第2页：目录 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title(slide, "今天聊什么？")
add_decorative_line(slide, Inches(1.4))
add_bullets(slide, [
    "什么是 Vibe Coding？它为何改变开发方式？",
    "大模型、Agent、RAG 等热词到底是什么？",
    "以 CCDC AI Service 为蓝本，看一个 AI 平台如何从零搭建",
    "Bond Codin：一个垂直领域 Agent 的设计实践",
    "踩坑经验、技术选型与未来展望"
])
add_footer(slide)

# ========== 第3页：什么是 Vibe Coding ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title(slide, "什么是 Vibe Coding？", "Describe, don't type.")
add_decorative_line(slide, Inches(1.4))
add_bullets(slide, [
    "Vibe Coding = 用自然语言描述意图，AI 辅助生成可运行代码",
    "开发者从『逐行手写代码』转向『描述需求 + 审核代码』",
    "核心三要素：清晰的 Prompt、迭代式对话、人的架构把控",
    "不是完全替代程序员，而是把重复性编码交给 AI，让人专注设计",
    "适用场景：原型搭建、CRUD、模板化代码、快速验证想法"
])
add_footer(slide)

# ========== 第4页：传统开发 vs Vibe Coding ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title(slide, "传统开发 vs Vibe Coding")
add_decorative_line(slide, Inches(1.4))

# 左侧：传统开发
left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.7), Inches(5.5), Inches(5.0))
tf = left_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "传统开发"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = TEXT_SUB
p.font.name = "Microsoft YaHei"
for text in ["需求 → 设计 → 手写代码 → 调试 → 上线", "开发者是主要代码生产者", "周期长，细节多，容易卡在小问题", "适合复杂核心业务、强定制化场景"]:
    p = tf.add_paragraph()
    p.text = "• " + text
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(10)

# 右侧：Vibe Coding
right_box = slide.shapes.add_textbox(Inches(7.0), Inches(1.7), Inches(5.5), Inches(5.0))
tf = right_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Vibe Coding"
p.font.size = Pt(22)
p.font.bold = True
p.font.color.rgb = ACCENT_CYAN
p.font.name = "Microsoft YaHei"
for text in ["意图 → AI 生成代码 → 人审核调优 → 上线", "开发者是架构师 + 审查员", "速度快，适合快速验证和 MVP", "适合工具类、原型、标准化模块"]:
    p = tf.add_paragraph()
    p.text = "• " + text
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_MAIN
    p.font.name = "Microsoft YaHei"
    p.space_after = Pt(10)

# 中间箭头
arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(6.0), Inches(3.5), Inches(1.0), Inches(0.5))
arrow.fill.solid()
arrow.fill.fore_color.rgb = ACCENT_PURPLE
arrow.line.fill.background()

add_footer(slide)

# ========== 第5页：大模型时代的新开发范式 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title(slide, "大模型时代的新开发范式")
add_decorative_line(slide, Inches(1.4))
add_bullets(slide, [
    "LLM（大语言模型）让机器真正『理解』自然语言并生成结构化输出",
    "开发范式演进：命令式编程 → 声明式编程 → 意图式编程",
    "Prompt 成为新的『编程语言』，Prompt Engineering 成为核心技能",
    "Agent = LLM + 工具调用 + 记忆 + 规划，能自主完成多步骤任务",
    "RAG（检索增强生成）让大模型『外挂』企业知识，解决幻觉问题"
])
add_footer(slide)

# ========== 第6页：项目背景 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title(slide, "项目背景：CCDC AI Service")
add_decorative_line(slide, Inches(1.4))
add_bullets(slide, [
    "项目代号：小发 —— 中债智能金融助手",
    "目标用户：债券从业者、发行承销商、固收分析师",
    "核心场景：债券定价预测、信用分析、募集书解读、宏观点评、收益率曲线",
    "技术目标：多用户在线、对话历史持久化、知识库 RAG、专家 Agent",
    "部署目标：Render 免费版，快速上线、低成本运维"
])
add_footer(slide)

# ========== 第7页：技术栈与架构 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title(slide, "技术栈与架构")
add_decorative_line(slide, Inches(1.4))
add_bullets(slide, [
    "前端：纯 HTML/CSS/JS 单页应用（无框架，轻量可控）",
    "后端：Node.js + Express（REST API + JWT 认证）",
    "数据库：PostgreSQL（Render 免费版，持久化存储）",
    "AI 接口：OpenAI 兼容协议，默认 Kimi + DeepSeek 双模型",
    "部署：Render.com（自动 CI/CD，免费实例休眠唤醒）",
    "关键依赖：pg、jsonwebtoken、bcryptjs、express"
])
add_footer(slide)

# ========== 第8页：零代码手写的真相 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title(slide, "零代码手写？真相是……", "Human in the Loop")
add_decorative_line(slide, Inches(1.4))
add_bullets(slide, [
    "不是『不写代码』，而是『不写重复代码』",
    "AI 负责生成：路由、CRUD、表单、样式骨架、Prompt 模板",
    "人类负责把控：架构设计、数据安全、业务规则、异常处理",
    "本项目实践：通过对话迭代生成代码 → 本地验证 → 部署上线",
    "核心能力：把业务意图准确翻译成 Prompt，并持续调试"
])
add_footer(slide)

# ========== 第9页：核心功能演进时间线 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title(slide, "核心功能演进时间线")
add_decorative_line(slide, Inches(1.4))
add_bullets(slide, [
    "V1：基础 AI 对话（标准 / 思考 / 预测三种模式）",
    "V2：用户注册登录 + JWT + 多用户数据隔离",
    "V3：PostgreSQL 持久化（解决 Render 休眠丢数据问题）",
    "V4：知识库 RAG + 多格式文件上传（docx/pdf/图片 OCR）",
    "V5：多模型切换（Kimi K2.6 + DeepSeek V4 Pro）",
    "V6：Bond Codin 专家 Agent —— 中债系统编程大师"
])
add_footer(slide)

# ========== 第10页：Bond Codin Agent 设计 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title(slide, "Agent 设计：Bond Codin", "中债系统编程大师")
add_decorative_line(slide, Inches(1.4))
add_bullets(slide, [
    "Agent 定义：具备角色、记忆、工具、目标的 AI 实体",
    "Bond Codin 融合三重角色：架构师 + 代码审计专家 + 分布式顾问",
    "核心能力：SIS-APP 三层架构分析、代码漏洞扫描、分布式改造建议",
    "交互设计：独立面板、代码编辑器、语言选择、文件上传",
    "输出规范：风险等级标注 + 四段式审计 + 改造蓝图"
])
add_footer(slide)

# ========== 第11页：Prompt Engineering 实践 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title(slide, "Prompt Engineering 实践")
add_decorative_line(slide, Inches(1.4))
add_bullets(slide, [
    "角色定义：你是 Bond Codin，一位拥有 20 年经验的金融 IT 架构师",
    "知识体系：Service/Logic/Mapper、MyBatis、Spring、分布式、云原生",
    "输出约束：必须标注 🔴严重 / 🟠高危 / 🟡中危 / 🟢建议",
    "结构化输出：问题定位 → 根因剖析 → 修复方案 → 架构影响",
    "少样本提示：通过示例让模型理解中债业务语境"
])
add_footer(slide)

# ========== 第12页：知识库与 RAG ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title(slide, "知识库与 RAG：让 AI 读懂企业文档")
add_decorative_line(slide, Inches(1.4))
add_bullets(slide, [
    "RAG = 检索（Retrieval）+ 生成（Generation）",
    "用户上传 docx/pdf/图片等文档，提取文本后存入数据库",
    "提问时做关键词匹配，把相关文本片段拼接到 Prompt",
    "本项目优化：中文按字拆分匹配，英文按词匹配，提升召回率",
    "价值：减少模型幻觉，让回答基于企业真实资料"
])
add_footer(slide)

# ========== 第13页：踩坑与经验 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title(slide, "踩坑与经验")
add_decorative_line(slide, Inches(1.4))
add_bullets(slide, [
    "坑 1：Render 免费版文件系统不持久 → 迁移到 PostgreSQL",
    "坑 2：JWT Secret 每次部署变化 → 固定环境变量",
    "坑 3：中文检索按空格分词失效 → 改为按字拆分匹配",
    "坑 4：某些模型只支持 temperature=1 → 统一设置为 1",
    "坑 5：Bond Codin 历史会话状态管理 → 统一 chatHistory + mode 字段",
    "经验：AI 生成代码后要人工验证边界条件和异常处理"
])
add_footer(slide)

# ========== 第14页：未来展望 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title(slide, "未来展望")
add_decorative_line(slide, Inches(1.4))
add_bullets(slide, [
    "MCP（Model Context Protocol）：标准化模型与工具的交互协议",
    "多 Agent 协作：多个专业 Agent 分工完成复杂任务",
    "本地模型 + 云端混合部署：敏感数据本地化，通用能力上云",
    "与企业知识库深度融合：实时同步内部文档、制度、规范",
    "AI 辅助编码进一步融入日常：从 Copilot 到全自动工作流"
])
add_footer(slide)

# ========== 第15页：总结 ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_title(slide, "总结")
add_decorative_line(slide, Inches(1.4))
add_bullets(slide, [
    "Vibe Coding 不是魔法，而是『意图驱动 + AI 生成 + 人审核』的新协作模式",
    "大模型让个人开发者也能快速搭建复杂应用原型",
    "Agent 设计的关键：明确的角色、清晰的输入输出、可控的边界",
    "企业级落地仍需关注：数据持久化、安全认证、模型边界、幻觉控制",
    "最重要的是：保持学习，把 AI 当作放大器，而不是替代品"
])
add_footer(slide)

# ========== 第16页：Q&A ==========
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
# 装饰
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.0), Inches(2.0), Inches(3.5), Inches(3.5))
circle.fill.solid()
circle.fill.fore_color.rgb = ACCENT_CYAN
circle.fill.transparency = 0.9
circle.line.color.rgb = ACCENT_CYAN
circle.line.width = Pt(2)

box = slide.shapes.add_textbox(Inches(0.8), Inches(2.8), Inches(11.7), Inches(2.0))
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Q & A"
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = "Microsoft YaHei"
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "欢迎交流讨论"
p2.font.size = Pt(24)
p2.font.color.rgb = ACCENT_CYAN
p2.font.name = "Microsoft YaHei"
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(16)

add_footer(slide)

# 保存
output_path = "CCDC_AI_VibeCoding_Training.pptx"
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
print(f"总页数：{len(prs.slides)}")
