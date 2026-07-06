# -*- coding: utf-8 -*-
"""
优化 PPT 视觉效果：
- 修复字体与背景撞色
- 统一全篇字体
- 提升对比度与可读性
- 增强卡片立体感
- 微调布局
不修改任何文字内容。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml
import re

INPUT_PATH = "CCDC_AI_VibeCoding_Training_v3.pptx"
OUTPUT_PATH = "CCDC_AI_VibeCoding_Training_v4.pptx"

prs = Presentation(INPUT_PATH)
W = prs.slide_width
H = prs.slide_height

# 优化后的配色
WHITE = RGBColor(0xff, 0xff, 0xff)
TEXT_MAIN = RGBColor(0xe6, 0xf7, 0xff)     # 主文字：浅青白
TEXT_BRIGHT = RGBColor(0xb8, 0xd4, 0xec)   # 次要文字：提亮后的灰蓝
CYAN = RGBColor(0x00, 0xe5, 0xff)          # 强调青
PURPLE = RGBColor(0xa6, 0x6f, 0xff)        # 强调紫（提亮）
ORANGE = RGBColor(0xff, 0xa8, 0x33)        # 强调橙（提亮）
GREEN = RGBColor(0x4d, 0xf0, 0x8a)         # 强调绿（提亮）
RED = RGBColor(0xff, 0x6b, 0x6b)           # 强调红（提亮）
CODE_TEXT = RGBColor(0xc8, 0xf8, 0xff)     # 代码文字

# 字体
FONT_CN = "Microsoft YaHei"
FONT_EN = "Segoe UI"
FONT_CODE = "Consolas"

# 需要修复的颜色映射（深色→浅色）
DARK_COLORS = {
    "060B1F": TEXT_MAIN,
    "000000": TEXT_MAIN,
    "1C2E52": TEXT_BRIGHT,
    "24355C": TEXT_BRIGHT,
    "8B9EBA": TEXT_BRIGHT,  # 提亮次要文字
}

# 强调色提亮映射
ACCENT_BRIGHTEN = {
    "8B5CFF": PURPLE,
    "FF9500": ORANGE,
    "00D26A": GREEN,
    "FF4D4F": RED,
    "A5F3FC": CODE_TEXT,
}


def is_code_text(text):
    """判断文本是否为代码（包含代码字符）"""
    if not text:
        return False
    code_chars = set("{}[]();=<>/\\|&+-*\"'#@")
    return any(c in text for c in code_chars) or text.count(".") >= 2


def fix_color(rgb_str):
    """修复与背景撞色的颜色"""
    s = str(rgb_str).upper()
    if s in DARK_COLORS:
        return DARK_COLORS[s]
    if s in ACCENT_BRIGHTEN:
        return ACCENT_BRIGHTEN[s]
    return None


def set_font(run, text):
    """统一字体"""
    if is_code_text(text):
        run.font.name = FONT_CODE
    else:
        # 简单判断：如果文本包含中文字符，用中文字体；否则用英文
        if re.search(r'[\u4e00-\u9fa5]', text):
            run.font.name = FONT_CN
        else:
            run.font.name = FONT_EN
        # 同时设置中文字体属性（对 Office 生效）
        try:
            run._element.set(qn('a:typeface'), FONT_CN)
        except Exception:
            pass


def optimize_text_frame(text_frame):
    """优化文本框内的字体、颜色、段落格式"""
    text_frame.word_wrap = True
    # 增加文本框内边距，避免文字贴边
    text_frame.margin_left = Inches(0.06)
    text_frame.margin_right = Inches(0.06)
    text_frame.margin_top = Inches(0.04)
    text_frame.margin_bottom = Inches(0.04)
    for para in text_frame.paragraphs:
        # 优化段落级字体
        if para.font.size and para.font.size < Pt(10):
            para.font.size = Pt(10)
        if para.font.size and para.font.size > Pt(44):
            para.font.size = Pt(44)

        # 修复段落级颜色
        try:
            if para.font.color.type is not None and para.font.color.rgb:
                fixed = fix_color(para.font.color.rgb)
                if fixed:
                    para.font.color.rgb = fixed
        except Exception:
            pass

        # 设置段落级字体
        if para.text.strip():
            try:
                if is_code_text(para.text):
                    para.font.name = FONT_CODE
                else:
                    para.font.name = FONT_CN
            except Exception:
                pass

        # 行距和段间距
        para.space_before = Pt(3)
        para.space_after = Pt(8)
        para.line_spacing = 1.25

        for run in para.runs:
            text = run.text
            if not text:
                continue

            # 统一字体
            set_font(run, text)

            # 修复颜色
            try:
                if run.font.color.type is not None and run.font.color.rgb:
                    fixed = fix_color(run.font.color.rgb)
                    if fixed:
                        run.font.color.rgb = fixed
            except Exception:
                pass

            # 确保代码不要太小
            if is_code_text(text) and run.font.size and run.font.size < Pt(9):
                run.font.size = Pt(9)

            # 粗体/斜体保持
            if run.font.bold is None and is_heading_text(text):
                run.font.bold = True


def is_heading_text(text):
    """简单判断是否为标题文字"""
    if not text:
        return False
    # 包含中文且无标点，长度较短，可能是标题
    if re.search(r'[\u4e00-\u9fa5]', text) and len(text) < 25 and not re.search(r'[，。；：]', text):
        return True
    return False


def is_card_shape(shape):
    """判断是否为卡片形状（圆角矩形/矩形、有填充、非背景装饰）"""
    try:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            return False
        # 排除全页背景、top bar、网格线、光晕圆、小线条
        if shape.width > W * 0.95 and shape.height > H * 0.95:
            return False
        if shape.height < Inches(0.2):
            return False
        if shape.width < Inches(0.3):
            return False
        # 只处理矩形/圆角矩形类
        shape_name = str(shape.auto_shape_type).lower() if hasattr(shape, 'auto_shape_type') else ''
        if 'rect' not in shape_name and 'round' not in shape_name and shape_name:
            return False
        # 有填充的深色卡片
        if shape.fill.type == 1:  # solid
            return True
        if shape.fill.type == 3:  # gradient
            return True
        return False
    except Exception:
        return False


def add_card_depth(slide, shape):
    """为卡片添加立体感：增强边框、可选渐变"""
    try:
        # 增强边框
        shape.line.color.rgb = RGBColor(0x3a, 0x55, 0x8a)
        shape.line.width = Pt(1.25)

        # 为卡片添加轻微渐变立体感
        if shape.fill.type == 1:  # solid
            rgb = str(shape.fill.fore_color.rgb).upper()
            dark_card_colors = ['0F172E', '141E3A', '1A3A5A', '1A4A5A', '1A5A5A', '1A5A4A',
                               '2A1515', '102A15', '1A3A6A', '1A4A6A', '1A4A4A',
                               '061020', '0A1228', '0F172E', '1A2A4A']
            if rgb in dark_card_colors:
                shape.fill.gradient()
                shape.fill.gradient_angle = 135
                stops = list(shape.fill.gradient_stops)
                if len(stops) >= 2:
                    stops[0].position = 0.0
                    stops[0].color.rgb = RGBColor(0x1c, 0x2c, 0x50)
                    stops[1].position = 1.0
                    stops[1].color.rgb = RGBColor(0x10, 0x18, 0x32)
        # 无填充卡片加微弱填充
        elif shape.fill.type is None or shape.fill.type == 5:
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(0x12, 0x1a, 0x32)
            shape.fill.transparency = 0.05
    except Exception:
        pass


def is_title_shape(shape):
    """判断是否为页面标题文本框"""
    try:
        if not shape.has_text_frame:
            return False
        text = shape.text_frame.text.strip()
        if not text:
            return False
        # 标题通常在页面顶部，字号大
        if shape.top < Inches(1.2):
            for para in shape.text_frame.paragraphs:
                if para.font.size and para.font.size >= Pt(20):
                    return True
        return False
    except Exception:
        return False


def is_footer_shape(shape):
    """判断是否为页脚"""
    try:
        if not shape.has_text_frame:
            return False
        return shape.top > Inches(7.0)
    except Exception:
        return False


# 处理每一页
for slide_idx, slide in enumerate(prs.slides, 1):
    print(f"处理第 {slide_idx} 页...")

    for shape in slide.shapes:
        # 跳过非文本形状但处理卡片
        if shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            if is_card_shape(shape):
                add_card_depth(slide, shape)

        # 优化文本
        if not shape.has_text_frame:
            continue

        tf = shape.text_frame
        text = tf.text.strip()
        if not text:
            continue

        optimize_text_frame(tf)

        # 标题增强
        if is_title_shape(shape):
            for para in tf.paragraphs:
                para.alignment = PP_ALIGN.LEFT
                if para.font.size and para.font.size >= Pt(28):
                    para.font.bold = True
                    para.font.color.rgb = WHITE
                elif para.font.size and para.font.size >= Pt(14):
                    para.font.color.rgb = TEXT_BRIGHT
                para.space_after = Pt(8)

        # 页脚字号保持较小
        if is_footer_shape(shape):
            for para in tf.paragraphs:
                if para.font.size and para.font.size > Pt(11):
                    para.font.size = Pt(11)
                para.font.color.rgb = TEXT_BRIGHT

prs.save(OUTPUT_PATH)
print(f"\n优化完成：{OUTPUT_PATH}")
print(f"总页数：{len(prs.slides)}")
