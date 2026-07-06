# -*- coding: utf-8 -*-
"""
CCDC AI Service · Vibe Coding 技术培训 PPT 生成脚本（第二版）
目标：内容更深入、结合项目实例、UI 更专业美观。
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ============================================================================
# 基础配置
# ============================================================================
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# 科技深蓝配色
BG = RGBColor(0x06, 0x0b, 0x1f)         # 背景
CARD_BG = RGBColor(0x0f, 0x17, 0x2e)    # 卡片背景
CARD_BG2 = RGBColor(0x14, 0x1e, 0x3a)   # 卡片背景2
CYAN = RGBColor(0x00, 0xe5, 0xff)       # 强调青
PURPLE = RGBColor(0x8b, 0x5c, 0xff)     # 强调紫
ORANGE = RGBColor(0xff, 0x95, 0x00)     # 强调橙
GREEN = RGBColor(0x00, 0xd2, 0x6a)      # 成功绿
RED = RGBColor(0xff, 0x4d, 0x4f)        # 告警红
TEXT_MAIN = RGBColor(0xe6, 0xf7, 0xff)  # 主文字
TEXT_SUB = RGBColor(0x8b, 0x9e, 0xba)   # 次要文字
TEXT_CODE = RGBColor(0xa5, 0xf3, 0xfc)  # 代码文字
WHITE = RGBColor(0xff, 0xff, 0xff)

# 字号
TITLE_SIZE = Pt(34)
SUBTITLE_SIZE = Pt(16)
BODY_SIZE = Pt(15)
SMALL_SIZE = Pt(13)
CODE_SIZE = Pt(11)
SECTION_SIZE = Pt(13)

# 字体
FONT = "Microsoft YaHei"
FONT_EN = "Segoe UI"

# 通用位置
MARGIN_L = Inches(0.65)
MARGIN_R = prs.slide_width - Inches(0.65)
CONTENT_TOP = Inches(1.35)

# ============================================================================
# 工具函数
# ============================================================================
def move_to_bottom(shape):
    """将形状移到最底层"""
    spTree = shape.part._element
    # 需要从 slide.shapes._spTree 中移除并插入
    sp = shape._element
    spTree = shape._element.getparent()
    spTree.remove(sp)
    spTree.insert(2, sp)


def add_background(slide, color=BG):
    """添加纯色背景"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()
    move_to_bottom(bg)
    return bg


def add_top_bar(slide, accent=CYAN):
    """顶部装饰条"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    move_to_bottom(bar)


def add_glow_circle(slide, x, y, w, h, color, transparency=0.92):
    """添加发光/装饰圆"""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.fill.transparency = transparency
    c.line.color.rgb = color
    c.line.width = Pt(1)
    c.line.color.brightness = 0.2
    move_to_bottom(c)
    return c


def add_card(slide, x, y, w, h, color=CARD_BG):
    """添加卡片背景"""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.color.rgb = RGBColor(0x24, 0x35, 0x5c)
    card.line.width = Pt(1)
    card.adjustments[0] = 0.05
    move_to_bottom(card)
    return card


def add_section_tag(slide, text, x, y, color=CYAN):
    """添加小标签"""
    tag = slide.shapes.add_textbox(x, y, Inches(3.0), Inches(0.35))
    tf = tag.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = SECTION_SIZE
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = FONT
    return tag


def add_title(slide, title, subtitle=None, accent=CYAN):
    """添加页面标题"""
    add_top_bar(slide, accent)
    box = slide.shapes.add_textbox(MARGIN_L, Inches(0.45), Inches(12.0), Inches(0.8))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = TITLE_SIZE
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = SUBTITLE_SIZE
        p2.font.color.rgb = TEXT_SUB
        p2.font.name = FONT
        p2.space_before = Pt(6)


def add_footer(slide, page_no, text="CCDC AI Service · Vibe Coding 实战"):
    """添加页脚与页码"""
    # 分隔线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, Inches(7.15), Inches(12.0), Inches(0.008))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x24, 0x35, 0x5c)
    line.line.fill.background()
    # 左侧文字
    box = slide.shapes.add_textbox(MARGIN_L, Inches(7.18), Inches(8.0), Inches(0.25))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_SUB
    p.font.name = FONT
    # 页码
    box2 = slide.shapes.add_textbox(Inches(12.0), Inches(7.18), Inches(0.6), Inches(0.25))
    tf2 = box2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = str(page_no)
    p2.font.size = Pt(11)
    p2.font.bold = True
    p2.font.color.rgb = CYAN
    p2.font.name = FONT
    p2.alignment = PP_ALIGN.RIGHT


def add_bullets(slide, bullets, x, y, w, h, font_size=BODY_SIZE, color=TEXT_MAIN, line_space=12):
    """添加列表文本"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, text in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = font_size
        p.font.color.rgb = color
        p.font.name = FONT
        p.space_after = Pt(line_space)
    return box


def add_code_block(slide, code, x, y, w, h):
    """添加代码块样式文本"""
    card = add_card(slide, x, y, w, h, RGBColor(0x08, 0x10, 0x22))
    # 覆盖 card 的边框，代码块用更暗的边框
    for sp in slide.shapes:
        if sp == card:
            sp.line.color.rgb = RGBColor(0x1a, 0x2a, 0x4a)
    box = slide.shapes.add_textbox(x + Inches(0.15), y + Inches(0.12), w - Inches(0.3), h - Inches(0.24))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = CODE_SIZE
    p.font.color.rgb = TEXT_CODE
    p.font.name = "Consolas"
    p.line_spacing = 1.15
    return box, card


def add_two_col_text(slide, left_title, left_items, right_title, right_items, y_top=CONTENT_TOP, accent_left=CYAN, accent_right=PURPLE):
    """左右双栏对比"""
    # 左栏
    add_section_tag(slide, left_title, MARGIN_L, y_top, accent_left)
    add_bullets(slide, left_items, MARGIN_L, y_top + Inches(0.38), Inches(5.8), Inches(5.0))
    # 右栏
    add_section_tag(slide, right_title, Inches(7.0), y_top, accent_right)
    add_bullets(slide, right_items, Inches(7.0), y_top + Inches(0.38), Inches(5.8), Inches(5.0))


def add_page(slide_index, title, subtitle=None, accent=CYAN):
    """创建一页并返回 slide，带统一装饰"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_glow_circle(slide, Inches(10.5), Inches(-1.2), Inches(4.5), Inches(4.5), accent, 0.94)
    add_title(slide, title, subtitle, accent)
    add_footer(slide, slide_index)
    return slide


# ============================================================================
# 第1页：封面
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_top_bar(slide, CYAN)
add_glow_circle(slide, Inches(9.2), Inches(1.0), Inches(4.2), Inches(4.2), PURPLE, 0.88)
add_glow_circle(slide, Inches(-1.5), Inches(4.5), Inches(3.5), Inches(3.5), CYAN, 0.92)

# 主标题区域
title_box = slide.shapes.add_textbox(MARGIN_L, Inches(2.1), Inches(11.5), Inches(2.2))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Vibe Coding 实战"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = FONT
p2 = tf.add_paragraph()
p2.text = "零代码手写搭建 CCDC AI 服务平台"
p2.font.size = Pt(32)
p2.font.color.rgb = CYAN
p2.font.name = FONT
p2.space_before = Pt(10)
p3 = tf.add_paragraph()
p3.text = "从想法到上线的 AI 应用开发新范式"
p3.font.size = Pt(18)
p3.font.color.rgb = TEXT_SUB
p3.font.name = FONT
p3.space_before = Pt(16)

# 副信息
info_box = slide.shapes.add_textbox(MARGIN_L, Inches(5.3), Inches(8.0), Inches(0.8))
tf = info_box.text_frame
p = tf.paragraphs[0]
p.text = "以 ccdc-ai-service 项目为蓝本"
p.font.size = Pt(14)
p.font.color.rgb = TEXT_SUB
p.font.name = FONT
p2 = tf.add_paragraph()
p2.text = "面向有经验开发者的技术分享"
p2.font.size = Pt(14)
p2.font.color.rgb = TEXT_SUB
p2.font.name = FONT
p2.space_before = Pt(4)

add_footer(slide, 1)

# ============================================================================
# 第2页：目录
# ============================================================================
slide = add_page(2, "议程", "三个篇章，从理念到落地")

sections = [
    ("01", "理念篇", "Vibe Coding 是什么、为什么、怎么做", CYAN),
    ("02", "项目篇", "ccdc-ai-service 的架构、实现与演进", PURPLE),
    ("03", "实战篇", "踩坑、经验、建议与未来展望", ORANGE),
]

for i, (num, sec_title, desc, color) in enumerate(sections):
    x = MARGIN_L + i * Inches(4.2)
    y = CONTENT_TOP + Inches(0.3)
    # 卡片
    card = add_card(slide, x, y, Inches(3.9), Inches(3.8), CARD_BG2)
    # 序号
    nb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(0.25), Inches(1.0), Inches(0.6))
    tf = nb.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(42)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = FONT_EN
    # 标题
    tb = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(1.1), Inches(3.4), Inches(0.5))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = sec_title
    p.font.size = Pt(22)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT
    # 描述
    db = slide.shapes.add_textbox(x + Inches(0.25), y + Inches(1.7), Inches(3.4), Inches(1.6))
    tf = db.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = desc
    p.font.size = BODY_SIZE
    p.font.color.rgb = TEXT_SUB
    p.font.name = FONT

# 底部补充
add_bullets(slide, [
    "目标听众：有开发经验的技术骨干 · 时长：约 45 分钟 · 形式：理论 + 项目代码解析"
], MARGIN_L, Inches(6.05), Inches(12.0), Inches(0.5), SMALL_SIZE, TEXT_SUB)


# ============================================================================
# 第3页：为什么讲 Vibe Coding？
# ============================================================================
slide = add_page(3, "为什么讲 Vibe Coding？", "传统软件开发正在发生的范式转移", CYAN)

add_bullets(slide, [
    "AI 大模型的代码生成能力已达到『可用』水平：从片段补全到整文件生成",
    "有经验开发者的价值正在迁移：从『写每一行代码』到『定义问题 + 设计架构 + 审查 AI 产出』",
    "企业内部工具、原型验证、标准化模块的开发周期可以大幅缩短",
    "ccdc-ai-service 就是一个典型案例：一个人用自然语言驱动，完成了前端、后端、数据库、AI 集成、Agent 设计"
], MARGIN_L, CONTENT_TOP, Inches(12.0), Inches(2.2))

# 痛点卡片
pain_points = [
    ("重复性编码", "CRUD、表单校验、路由模板占据大量时间，创造力被消耗在低价值劳动中"),
    ("跨栈维护难", "前端 + 后端 + 数据库 + 部署，一人全栈时细节爆炸"),
    ("知识壁垒", "业务专家懂规则但不懂代码，技术专家懂代码但不懂金融业务"),
    ("试错成本高", "传统模式下，一个想法从立项到 Demo 往往需要数周"),
]
for i, (ptitle, pdesc) in enumerate(pain_points):
    x = MARGIN_L + (i % 2) * Inches(6.25)
    y = CONTENT_TOP + Inches(2.5) + (i // 2) * Inches(1.65)
    add_card(slide, x, y, Inches(6.0), Inches(1.45), CARD_BG2)
    tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(5.6), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = ptitle
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RED if i % 2 == 0 else ORANGE
    p.font.name = FONT
    tb2 = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.55), Inches(5.6), Inches(0.8))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = pdesc
    p2.font.size = SMALL_SIZE
    p2.font.color.rgb = TEXT_SUB
    p2.font.name = FONT

# ============================================================================
# 第4页：什么是 Vibe Coding？
# ============================================================================
slide = add_page(4, "什么是 Vibe Coding？", "Describe, don't type.", CYAN)

add_bullets(slide, [
    "Vibe Coding 指开发者用自然语言描述意图，由 AI 生成代码，开发者负责审查、调优与架构把控",
    "它不是『不写代码』，而是把编码工作从『逐字符敲击』转变为『意图表达 + 质量审核』",
    "成功的关键在于：清晰的上下文、可验证的目标、持续的反馈循环",
    "适合：原型搭建、标准模块、工具类应用；不适合：强安全关键系统、未成熟的业务逻辑"
], MARGIN_L, CONTENT_TOP, Inches(7.0), Inches(2.5))

# 右侧流程图
flow_steps = [
    ("1", "描述意图", "用自然语言说明想要什么"),
    ("2", "AI 生成", "模型输出代码/配置/测试"),
    ("3", "人审查", "检查正确性、安全性、可维护性"),
    ("4", "迭代调优", "补充上下文，修正偏差"),
    ("5", "集成上线", "嵌入项目，自动化验证"),
]
start_y = CONTENT_TOP
for i, (num, step, desc) in enumerate(flow_steps):
    y = start_y + i * Inches(0.85)
    # 序号圆
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(8.0), y, Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = CYAN if i % 2 == 0 else PURPLE
    circle.line.fill.background()
    tb = slide.shapes.add_textbox(Inches(8.0), y + Inches(0.05), Inches(0.45), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = BG
    p.font.name = FONT_EN
    p.alignment = PP_ALIGN.CENTER
    # 文字
    tb2 = slide.shapes.add_textbox(Inches(8.65), y, Inches(4.0), Inches(0.75))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = step
    p2.font.size = Pt(15)
    p2.font.bold = True
    p2.font.color.rgb = WHITE
    p2.font.name = FONT
    p3 = tf2.add_paragraph()
    p3.text = desc
    p3.font.size = SMALL_SIZE
    p3.font.color.rgb = TEXT_SUB
    p3.font.name = FONT
    p3.space_before = Pt(2)
    # 连接线
    if i < len(flow_steps) - 1:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.2), y + Inches(0.5), Inches(0.04), Inches(0.35))
        line.fill.solid()
        line.fill.fore_color.rgb = RGBColor(0x24, 0x35, 0x5c)
        line.line.fill.background()

# ============================================================================
# 第5页：开发范式演进
# ============================================================================
slide = add_page(5, "开发范式演进", "从命令式到意图式", CYAN)

timeline_y = CONTENT_TOP + Inches(0.6)
# 横线
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, timeline_y + Inches(0.35), Inches(12.0), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0x24, 0x35, 0x5c)
line.line.fill.background()

stages = [
    ("命令式", "Imperative", "逐行告诉计算机怎么做\n汇编 / C / 早期脚本", TEXT_SUB),
    ("声明式", "Declarative", "描述目标状态，框架执行\nSQL / React / 配置化", TEXT_SUB),
    ("低代码", "Low-Code", "拖拽 + 配置 + 少量脚本\n快速表单、工作流", ORANGE),
    ("意图式", "Intent-Driven", "自然语言描述，AI 实现\nVibe Coding / Agent", CYAN),
]
for i, (cn, en, desc, color) in enumerate(stages):
    x = MARGIN_L + i * Inches(3.05)
    # 节点
    node = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, timeline_y, Inches(0.7), Inches(0.7))
    node.fill.solid()
    node.fill.fore_color.rgb = color
    node.line.fill.background()
    # 标题
    tb = slide.shapes.add_textbox(x - Inches(0.2), timeline_y + Inches(0.9), Inches(1.5), Inches(0.4))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = cn
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    tb2 = slide.shapes.add_textbox(x - Inches(0.2), timeline_y + Inches(1.3), Inches(1.5), Inches(0.35))
    tf2 = tb2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = en
    p2.font.size = Pt(11)
    p2.font.color.rgb = TEXT_SUB
    p2.font.name = FONT_EN
    p2.alignment = PP_ALIGN.CENTER
    # 描述
    db = slide.shapes.add_textbox(x - Inches(0.45), timeline_y + Inches(1.75), Inches(2.0), Inches(1.0))
    dtf = db.text_frame
    dtf.word_wrap = True
    dp = dtf.paragraphs[0]
    dp.text = desc
    dp.font.size = SMALL_SIZE
    dp.font.color.rgb = TEXT_SUB
    dp.font.name = FONT
    dp.alignment = PP_ALIGN.CENTER

# 底部说明
add_bullets(slide, [
    "Vibe Coding 不是凭空出现，而是软件开发抽象层级持续提升的必然结果",
    "每一层抽象都伴随着质疑，但也释放了大量生产力；关键在于知道边界在哪里"
], MARGIN_L, Inches(5.6), Inches(12.0), Inches(1.0), BODY_SIZE, TEXT_MAIN)

# ============================================================================
# 第6页：项目全景
# ============================================================================
slide = add_page(6, "项目全景：CCDC AI Service", "代号『小发』——中债智能金融助手", PURPLE)

# 左侧：业务场景
add_section_tag(slide, "业务场景", MARGIN_L, CONTENT_TOP, CYAN)
add_bullets(slide, [
    "目标用户：债券承销商、固收研究员、发行机构业务人员",
    "核心功能：债券定价预测、信用分析、募集说明书解读、宏观点评、收益率曲线分析",
    "产品形态：Web 单页应用 + REST API + 实时流式对话",
    "部署目标：Render 免费版，零成本快速上线"
], MARGIN_L, CONTENT_TOP + Inches(0.38), Inches(5.8), Inches(2.6))

# 右侧：关键指标
add_section_tag(slide, "项目关键指标", Inches(7.0), CONTENT_TOP, PURPLE)
metrics = [
    ("约 3700 行", "前端代码（纯 HTML/CSS/JS）"),
    ("~600 行", "后端 Express 主服务"),
    ("4 张表", "PostgreSQL：users/chats/configs/kb_docs"),
    ("2 模型", "Kimi K2.6 + DeepSeek V4 Pro"),
    ("5+ 格式", "知识库支持 docx/pdf/图片 OCR"),
]
for i, (val, label) in enumerate(metrics):
    y = CONTENT_TOP + Inches(0.45) + i * Inches(0.55)
    # 数值
    vb = slide.shapes.add_textbox(Inches(7.0), y, Inches(1.6), Inches(0.4))
    vf = vb.text_frame
    vp = vf.paragraphs[0]
    vp.text = val
    vp.font.size = Pt(15)
    vp.font.bold = True
    vp.font.color.rgb = CYAN
    vp.font.name = FONT_EN
    # 标签
    lb = slide.shapes.add_textbox(Inches(8.7), y + Inches(0.04), Inches(4.0), Inches(0.4))
    lf = lb.text_frame
    lp = lf.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(14)
    lp.font.color.rgb = TEXT_MAIN
    lp.font.name = FONT

# 底部引用
add_card(slide, MARGIN_L, Inches(5.6), Inches(12.0), Inches(1.0), CARD_BG2)
add_bullets(slide, [
    "本项目不是 toy project，而是承载了真实业务诉求的完整 AI 应用：多用户、持久化、知识库、专家 Agent 一应俱全。"
], MARGIN_L + Inches(0.2), Inches(5.75), Inches(11.6), Inches(0.7), BODY_SIZE, TEXT_MAIN)

# ============================================================================
# 第7页：系统架构
# ============================================================================
slide = add_page(7, "系统架构：四层设计", "清晰分层是 AI 项目可维护的关键", PURPLE)

# 架构图：4 层
layers = [
    ("表现层", "public/index.html\n单页应用 · 主题切换 · 历史会话", RGBColor(0x1a, 0x3a, 0x5a)),
    ("接入层", "Express + JWT/CORS\nREST API · SSE 流式 · 认证中间件", RGBColor(0x1a, 0x4a, 0x5a)),
    ("业务层", "对话 · 知识库 · Bond Codin\n配置管理 · 文件解析", RGBColor(0x1a, 0x5a, 0x5a)),
    ("数据层", "PostgreSQL\nusers / chats / configs / kb_docs", RGBColor(0x1a, 0x5a, 0x4a)),
]
layer_w = Inches(10.5)
layer_h = Inches(0.95)
start_x = MARGIN_L + Inches(0.75)
for i, (name, desc, color) in enumerate(layers):
    y = CONTENT_TOP + i * Inches(1.15)
    # 层背景
    rect = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, start_x, y, layer_w, layer_h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = color
    rect.line.color.rgb = CYAN if i == 1 else RGBColor(0x2a, 0x4a, 0x6a)
    rect.line.width = Pt(1)
    # 名称
    nb = slide.shapes.add_textbox(start_x + Inches(0.25), y + Inches(0.12), Inches(1.6), Inches(0.4))
    nf = nb.text_frame
    np = nf.paragraphs[0]
    np.text = name
    np.font.size = Pt(16)
    np.font.bold = True
    np.font.color.rgb = WHITE
    np.font.name = FONT
    # 描述
    db = slide.shapes.add_textbox(start_x + Inches(2.0), y + Inches(0.22), Inches(8.2), Inches(0.5))
    df = db.text_frame
    dp = df.paragraphs[0]
    dp.text = desc
    dp.font.size = Pt(13)
    dp.font.color.rgb = TEXT_MAIN
    dp.font.name = FONT

# 左侧标签
for i, _ in enumerate(layers):
    y = CONTENT_TOP + i * Inches(1.15) + Inches(0.25)
    tb = slide.shapes.add_textbox(MARGIN_L, y, Inches(0.6), Inches(0.45))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = f"L{i+1}"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.font.name = FONT_EN

# 右侧外部依赖
add_card(slide, Inches(10.5), CONTENT_TOP + Inches(0.2), Inches(2.35), Inches(3.6), CARD_BG2)
add_bullets(slide, [
    "外部依赖",
    "",
    "• Moonshot API",
    "• DeepSeek API",
    "• Render.com",
    "• PostgreSQL"
], Inches(10.65), CONTENT_TOP + Inches(0.35), Inches(2.0), Inches(3.3), SMALL_SIZE, TEXT_SUB)


# ============================================================================
# 第8页：前端零框架设计
# ============================================================================
slide = add_page(8, "前端为何零框架？", "public/index.html 的取舍与坚持", PURPLE)

add_two_col_text(
    slide,
    "选择零框架的理由",
    [
        "单文件即可承载完整 SPA，部署极简",
        "无构建步骤，修改后立即刷新，迭代极快",
        "对 Vibe Coding 友好：AI 生成单文件 HTML 的成功率高",
        "避免 React/Vue 生态的版本冲突和构建配置"
    ],
    "带来的挑战",
    [
        "状态管理需手写，3700 行代码中业务逻辑密集",
        "组件复用靠函数封装，容易重复",
        "没有类型系统，重构风险高",
        "需要手动处理 XSS、事件委托等细节"
    ],
    CONTENT_TOP,
    GREEN,
    ORANGE
)

# 代码示例
code = """// 主题切换：纯 JS 操作 class 与 localStorage
function loadTheme() {
  const t = localStorage.getItem(THEME_KEY) || 'light';
  document.body.classList.toggle('dark', t === 'dark');
}"""
add_code_block(slide, code, MARGIN_L, Inches(4.85), Inches(12.0), Inches(1.7))

add_bullets(slide, [
    "结论：对于 MVP 和内部工具，零框架可以显著降低认知负担；当复杂度超过临界点时，再引入框架。"
], MARGIN_L, Inches(6.7), Inches(12.0), Inches(0.6), SMALL_SIZE, TEXT_SUB)

# ============================================================================
# 第9页：后端 Express 核心设计
# ============================================================================
slide = add_page(9, "后端 Express 核心设计", "轻量但完整的 API 层", PURPLE)

add_bullets(slide, [
    "路由分组：/api/auth、/api/chat、/api/history、/api/config、/api/kb",
    "认证中间件：JWT 校验 + 用户数据隔离",
    "统一错误处理：try/catch + 结构化错误响应",
    "流式响应：/api/chat 使用 SSE 输出大模型结果，降低首字时延"
], MARGIN_L, CONTENT_TOP, Inches(6.3), Inches(2.3))

# 路由卡片
routes = [
    ("POST /api/auth/register", "注册 + 确认密码校验 + bcrypt 哈希"),
    ("POST /api/auth/login", "登录 + JWT 签发"),
    ("POST /api/chat", "知识库检索 + SSE 流式代理"),
    ("GET /api/history", "查询当前用户对话历史"),
    ("POST /api/kb/upload", "docx/pdf/图片 OCR 入库"),
]
for i, (route, desc) in enumerate(routes):
    y = CONTENT_TOP + i * Inches(0.85)
    add_card(slide, Inches(7.0), y, Inches(5.65), Inches(0.72), CARD_BG2)
    rb = slide.shapes.add_textbox(Inches(7.15), y + Inches(0.1), Inches(3.0), Inches(0.3))
    rf = rb.text_frame
    rp = rf.paragraphs[0]
    rp.text = route
    rp.font.size = Pt(12)
    rp.font.bold = True
    rp.font.color.rgb = CYAN
    rp.font.name = "Consolas"
    db = slide.shapes.add_textbox(Inches(7.15), y + Inches(0.4), Inches(5.3), Inches(0.3))
    df = db.text_frame
    dp = df.paragraphs[0]
    dp.text = desc
    dp.font.size = Pt(11)
    dp.font.color.rgb = TEXT_SUB
    dp.font.name = FONT

# 代码块
code = """// 认证中间件：统一校验 JWT
function authMiddleware(req, res, next) {
  const token = req.headers.authorization?.split(' ')[1];
  if (!token) return res.status(401).json({ error: '未登录' });
  try {
    req.user = jwt.verify(token, JWT_SECRET);
    next();
  } catch { res.status(401).json({ error: '登录已过期' }); }
}"""
add_code_block(slide, code, MARGIN_L, Inches(4.5), Inches(12.0), Inches(2.0))

# ============================================================================
# 第10页：数据库迁移实战
# ============================================================================
slide = add_page(10, "数据库迁移实战", "SQLite → PostgreSQL 的决策与实现", PURPLE)

# 决策原因
add_section_tag(slide, "为什么迁移？", MARGIN_L, CONTENT_TOP, RED)
add_bullets(slide, [
    "Render 免费版 Web Service 没有 Persistent Disk",
    "实例休眠后本地 SQLite 文件丢失，用户注册、对话历史全部归零",
    "PostgreSQL 作为独立托管服务，数据持久化且跨实例共享",
    "pg 模块与 Node.js 生态集成成熟，迁移成本低"
], MARGIN_L, CONTENT_TOP + Inches(0.38), Inches(12.0), Inches(1.5))

# 对比
add_two_col_text(
    slide,
    "SQLite 阶段",
    ["本地文件存储", "单进程访问", "零运维成本", "适合本地开发/测试"],
    "PostgreSQL 阶段",
    ["网络数据库服务", "多并发连接池", "Render 托管/自运维", "适合生产/多实例"],
    CONTENT_TOP + Inches(2.0),
    ORANGE,
    GREEN
)

# 代码示例
code = """const { Pool } = require('pg');
const pool = new Pool({
  connectionString: process.env.DATABASE_URL,
  ssl: process.env.DATABASE_URL?.includes('render.com')
       ? { rejectUnauthorized: false }
       : false
});"""
add_code_block(slide, code, MARGIN_L, Inches(5.1), Inches(12.0), Inches(1.35))

# ============================================================================
# 第11页：AI 层多模型抽象
# ============================================================================
slide = add_page(11, "AI 层多模型抽象", "如何优雅对接多个大模型", PURPLE)

add_bullets(slide, [
    "业务需求：用户可在 Kimi K2.6 与 DeepSeek V4 Pro 之间切换",
    "统一接口：后端兼容 OpenAI 风格的 /chat/completions 协议",
    "配置化：通过 ADMIN_API_MODELS 环境变量定义模型列表",
    "模型标识：kimi-k2.6|Kimi K2.6,deepseek-v4-pro|DeepSeek V4 Pro"
], MARGIN_L, CONTENT_TOP, Inches(12.0), Inches(1.8))

# 代码块
code = """function parseAdminModels() {
  const raw = process.env.ADMIN_API_MODELS;
  const models = raw ? raw.split(',') : [];
  const parsed = models.map(m => {
    const [id, name] = m.split('|');
    return { id: id.trim(), name: (name || id).trim() };
  });
  // 兼容旧版：自动追加 deepseek-v4-pro
  if (!parsed.find(m => m.id === 'deepseek-v4-pro')) {
    parsed.push({ id: 'deepseek-v4-pro', name: 'DeepSeek V4 Pro' });
  }
  return parsed;
}"""
add_code_block(slide, code, MARGIN_L, Inches(3.1), Inches(12.0), Inches(2.4))

add_bullets(slide, [
    "设计收益：新增模型只需改环境变量，无需改动业务代码",
    "注意点：不同模型对 temperature 等参数的支持存在差异，需统一为兼容值"
], MARGIN_L, Inches(5.75), Inches(12.0), Inches(0.8), BODY_SIZE, TEXT_MAIN)

# ============================================================================
# 第12页：SSE 流式响应
# ============================================================================
slide = add_page(12, "SSE 流式响应实现", "降低首字时延，提升对话体验", PURPLE)

add_bullets(slide, [
    "为什么用 SSE：大模型生成长文本时，用户无需等待完整响应",
    "前端 EventSource 逐段接收并渲染，呈现打字机效果",
    "后端将上游模型返回的流直接透传，不做全量缓冲",
    "异常处理：网络中断、模型超时、JSON 解析失败都需兜底"
], MARGIN_L, CONTENT_TOP, Inches(12.0), Inches(1.7))

# 时序图（简化版）
roles = [("前端", MARGIN_L + Inches(0.3)), ("Express", MARGIN_L + Inches(4.5)), ("Moonshot", MARGIN_L + Inches(8.7))]
for name, x in roles:
    # 角色标签
    tb = slide.shapes.add_textbox(x, Inches(3.45), Inches(1.2), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = name
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.font.name = FONT
    p.alignment = PP_ALIGN.CENTER
    # 生命线
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + Inches(0.5), Inches(3.8), Inches(0.02), Inches(2.3))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(0x30, 0x45, 0x6a)
    line.line.fill.background()

# 箭头
arrows = [
    (MARGIN_L + Inches(0.8), Inches(4.0), MARGIN_L + Inches(4.5), "POST /api/chat"),
    (MARGIN_L + Inches(5.0), Inches(4.4), MARGIN_L + Inches(8.7), "stream request"),
    (MARGIN_L + Inches(8.7), Inches(4.8), MARGIN_L + Inches(5.0), "chunk 1"),
    (MARGIN_L + Inches(4.5), Inches(5.2), MARGIN_L + Inches(0.8), "chunk 1"),
    (MARGIN_L + Inches(8.7), Inches(5.6), MARGIN_L + Inches(5.0), "chunk N"),
]
for x1, y, x2, label in arrows:
    is_forward = x2 > x1
    # 简化箭头用矩形代替
    # 用水平条表示数据流
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x1, y + Inches(0.02), abs(x2 - x1) - Inches(0.1), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = CYAN if is_forward else PURPLE
    bar.line.fill.background()
    # 标签
    lb = slide.shapes.add_textbox(min(x1, x2) + Inches(0.2), y - Inches(0.15), Inches(2.0), Inches(0.2))
    lf = lb.text_frame
    lp = lf.paragraphs[0]
    lp.text = label
    lp.font.size = Pt(9)
    lp.font.color.rgb = TEXT_SUB
    lp.font.name = FONT

# 注：由于 python-pptx 对箭头方向控制有限，SSE 流程用文字补充
add_bullets(slide, [
    "流程：前端请求 → Express 鉴权 → 转发 Moonshot/DeepSeek → 流式 chunks 逐段回传 → 前端实时渲染"
], MARGIN_L, Inches(6.2), Inches(12.0), Inches(0.5), SMALL_SIZE, TEXT_SUB)

# ============================================================================
# 第13页：知识库 RAG
# ============================================================================
slide = add_page(13, "知识库 RAG：让 AI 读懂企业文档", "检索增强生成的项目实践", PURPLE)

# 流程图：4 步
steps = [
    ("上传", "docx / pdf / 图片", "用户在前端选择文件"),
    ("解析", "文本/OCR 提取", "后端调用 docx/pdf/ocr 库"),
    ("存储", "PostgreSQL kb_docs", "原始文本 + 提取内容"),
    ("检索", "关键词匹配", "按查询词召回相关片段"),
]
for i, (step, detail, note) in enumerate(steps):
    x = MARGIN_L + i * Inches(3.1)
    y = CONTENT_TOP
    add_card(slide, x, y, Inches(2.85), Inches(2.4), CARD_BG2)
    # 步骤名
    sb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(2.45), Inches(0.4))
    sf = sb.text_frame
    sp = sf.paragraphs[0]
    sp.text = step
    sp.font.size = Pt(18)
    sp.font.bold = True
    sp.font.color.rgb = CYAN
    sp.font.name = FONT
    # 详情
    db = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), Inches(2.45), Inches(0.5))
    df = db.text_frame
    df.word_wrap = True
    dp = df.paragraphs[0]
    dp.text = detail
    dp.font.size = Pt(14)
    dp.font.bold = True
    dp.font.color.rgb = WHITE
    dp.font.name = FONT
    # 说明
    nb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(1.25), Inches(2.45), Inches(1.0))
    nf = nb.text_frame
    nf.word_wrap = True
    np = nf.paragraphs[0]
    np.text = note
    np.font.size = SMALL_SIZE
    np.font.color.rgb = TEXT_SUB
    np.font.name = FONT
    # 箭头
    if i < len(steps) - 1:
        arr = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.95), y + Inches(0.9), Inches(0.28), Inches(0.18))
        arr.fill.solid()
        arr.fill.fore_color.rgb = PURPLE
        arr.line.fill.background()

add_bullets(slide, [
    "RAG 的价值：把大模型从『通用知识』转变为『企业知识助手』，显著降低幻觉风险"
], MARGIN_L, Inches(4.4), Inches(12.0), Inches(0.5), BODY_SIZE, TEXT_MAIN)

# 代码块
code = """// 中文检索优化：按字拆分，英文按词拆分
function tokenize(text) {
  const zh = text.match(/[\u4e00-\u9fa5]/g) || [];
  const en = text.toLowerCase().match(/[a-z0-9]+/g) || [];
  return [...zh, ...en];
}"""
add_code_block(slide, code, MARGIN_L, Inches(4.95), Inches(12.0), Inches(1.5))

# ============================================================================
# 第14页：中文检索优化
# ============================================================================
slide = add_page(14, "中文检索优化", "分词策略直接影响 RAG 召回率", PURPLE)

add_two_col_text(
    slide,
    "优化前：按空格分词",
    [
        "将『债券发行定价』整体查询",
        "仅匹配包含完整短语的文档",
        "对 OCR 误差、同义词、口语化表达非常脆弱",
        "漏召率高，AI 经常『不知道』"
    ],
    "优化后：中文逐字 + 英文按词",
    [
        "查询拆成 『债/券/发/行/定/价』",
        "只要文档包含其中若干字即可命中",
        "召回率显著提升",
        "虽可能引入噪声，但对 RAG 场景利大于弊"
    ],
    CONTENT_TOP,
    RED,
    GREEN
)

# 示例对比
add_card(slide, MARGIN_L, Inches(4.6), Inches(5.9), Inches(2.0), RGBColor(0x2a, 0x15, 0x15))
add_bullets(slide, [
    "原查询：'债券发行定价'",
    "原分词：['债券发行定价']",
    "问题：文档写『债券的发行与定价』就无法命中"
], MARGIN_L + Inches(0.2), Inches(4.8), Inches(5.5), Inches(1.6), SMALL_SIZE, TEXT_SUB)

add_card(slide, Inches(6.9), Inches(4.6), Inches(5.9), Inches(2.0), RGBColor(0x10, 0x2a, 0x15))
add_bullets(slide, [
    "新查询：'债券发行定价'",
    "新分词：['债','券','发','行','定','价']",
    "效果：只要包含任意几个字，就会进入候选"
], Inches(7.1), Inches(4.8), Inches(5.5), Inches(1.6), SMALL_SIZE, TEXT_SUB)

# ============================================================================
# 第15页：Bond Codin Agent 设计
# ============================================================================
slide = add_page(15, "Agent 设计：Bond Codin", "中债系统编程大师", ORANGE)

add_bullets(slide, [
    "Agent 三要素：角色（Role）、记忆（Memory）、工具（Tools）",
    "Bond Codin 是面向中债科技业务领域的垂直 Agent",
    "它融合架构师、代码审计专家、分布式改造顾问三种角色",
    "核心使命：对中债系统代码进行架构分析、漏洞审计、改造建议"
], MARGIN_L, CONTENT_TOP, Inches(12.0), Inches(1.8))

# 角色卡片
roles = [
    ("架构师", "Service/Logic/Mapper 三层架构评审，识别职责混乱"),
    ("审计专家", "发现 SQL 注入、空指针、并发风险、越权访问"),
    ("改造顾问", "给出云原生拆分、缓存引入、异步化、分库分表建议"),
]
for i, (role, desc) in enumerate(roles):
    x = MARGIN_L + i * Inches(4.15)
    y = CONTENT_TOP + Inches(2.1)
    add_card(slide, x, y, Inches(3.95), Inches(2.0), CARD_BG2)
    rb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.2), Inches(3.55), Inches(0.4))
    rf = rb.text_frame
    rp = rf.paragraphs[0]
    rp.text = role
    rp.font.size = Pt(18)
    rp.font.bold = True
    rp.font.color.rgb = CYAN
    rp.font.name = FONT
    db = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.7), Inches(3.55), Inches(1.2))
    df = db.text_frame
    df.word_wrap = True
    dp = df.paragraphs[0]
    dp.text = desc
    dp.font.size = SMALL_SIZE
    dp.font.color.rgb = TEXT_MAIN
    dp.font.name = FONT

add_bullets(slide, [
    "入口设计：前端独立 Bond Codin 面板，支持代码编辑、语言选择、文件上传、历史会话"
], MARGIN_L, Inches(5.8), Inches(12.0), Inches(0.5), BODY_SIZE, TEXT_MAIN)


# ============================================================================
# 第16页：Prompt Engineering 深度解析
# ============================================================================
slide = add_page(16, "Prompt Engineering 深度解析", "Bond Codin 的提示词设计", ORANGE)

add_bullets(slide, [
    "Prompt 是 Vibe Coding 和 Agent 的核心交付物，决定了 AI 的行为边界",
    "Bond Codin 的 system prompt 采用『角色 + 知识 + 约束 + 输出格式』四段式",
    "通过环境变量注入，支持不修改代码即可调整专家行为"
], MARGIN_L, CONTENT_TOP, Inches(12.0), Inches(1.5))

# 提示词结构图
prompt_parts = [
    ("角色定义", "你是 Bond Codin，一位拥有 20 年经验的金融 IT 架构师……", CYAN),
    ("知识体系", "精通 SIS-APP 三层架构、MyBatis、Spring Cloud、K8s……", PURPLE),
    ("输出约束", "必须按风险等级标注：🔴严重 🟠高危 🟡中危 🟢建议", ORANGE),
    ("格式规范", "问题定位 → 根因剖析 → 修复方案 → 架构影响", GREEN),
]
for i, (part, example, color) in enumerate(prompt_parts):
    y = CONTENT_TOP + Inches(1.7) + i * Inches(0.95)
    add_card(slide, MARGIN_L, y, Inches(12.0), Inches(0.85), CARD_BG2)
    # 标签
    tb = slide.shapes.add_textbox(MARGIN_L + Inches(0.15), y + Inches(0.12), Inches(1.6), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = part
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = color
    p.font.name = FONT
    # 示例
    eb = slide.shapes.add_textbox(MARGIN_L + Inches(1.85), y + Inches(0.15), Inches(10.0), Inches(0.55))
    ef = eb.text_frame
    ef.word_wrap = True
    ep = ef.paragraphs[0]
    ep.text = example
    ep.font.size = Pt(12)
    ep.font.color.rgb = TEXT_MAIN
    ep.font.name = FONT

add_bullets(slide, [
    "技巧：用少量示例（Few-shot）让模型理解中债业务语境，减少泛化错误"
], MARGIN_L, Inches(6.65), Inches(12.0), Inches(0.5), SMALL_SIZE, TEXT_SUB)

# ============================================================================
# 第17页：代码审计输出规范
# ============================================================================
slide = add_page(17, "代码审计输出规范", "Bond Codin 如何给出可落地的结论", ORANGE)

# 四段式
sections_detail = [
    ("1. 问题定位", "指出代码行、函数、类，以及违反了哪条设计原则或规范", CYAN),
    ("2. 根因剖析", "解释为什么这是问题：并发模型、数据流、权限边界、性能瓶颈", PURPLE),
    ("3. 修复方案", "给出可直接修改的代码片段或重构步骤，标注风险等级", ORANGE),
    ("4. 架构影响", "评估对 Service/Logic/Mapper 三层、数据库、接口契约的影响", GREEN),
]
for i, (title, desc, color) in enumerate(sections_detail):
    y = CONTENT_TOP + i * Inches(1.1)
    add_card(slide, MARGIN_L, y, Inches(12.0), Inches(0.95), CARD_BG2)
    # 色块
    block = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L + Inches(0.12), y + Inches(0.15), Inches(0.12), Inches(0.65))
    block.fill.solid()
    block.fill.fore_color.rgb = color
    block.line.fill.background()
    # 标题
    tb = slide.shapes.add_textbox(MARGIN_L + Inches(0.35), y + Inches(0.12), Inches(2.4), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.font.name = FONT
    # 描述
    db = slide.shapes.add_textbox(MARGIN_L + Inches(2.85), y + Inches(0.15), Inches(9.0), Inches(0.65))
    df = db.text_frame
    df.word_wrap = True
    dp = df.paragraphs[0]
    dp.text = desc
    dp.font.size = SMALL_SIZE
    dp.font.color.rgb = TEXT_MAIN
    dp.font.name = FONT

add_bullets(slide, [
    "价值：让 AI 输出从『泛泛而谈』变成『可执行的技术方案』"
], MARGIN_L, Inches(6.75), Inches(12.0), Inches(0.5), BODY_SIZE, TEXT_MAIN)

# ============================================================================
# 第18页：文件上传与 OCR
# ============================================================================
slide = add_page(18, "文件上传与 OCR", "多格式知识库接入", ORANGE)

add_bullets(slide, [
    "业务诉求：用户上传募集说明书、审计报告、制度文件、代码截图等",
    "技术实现：根据 MIME 类型分发不同解析器",
    "docx：读取段落与表格文本",
    "pdf：文本层抽取，复杂版式可结合 OCR",
    "图片：OCR 提取后作为知识库文本存储"
], MARGIN_L, CONTENT_TOP, Inches(6.2), Inches(2.4))

# 文件类型卡片
file_types = [
    (".docx", "docx 解析库", "段落、表格、页眉页脚"),
    (".pdf", "pdf 解析库", "文本层、图片兜底"),
    (".png/.jpg", "OCR 服务", "代码截图、扫描件"),
    (".txt/.md", "直接读取", "制度、规范文档"),
]
for i, (ext, lib, desc) in enumerate(file_types):
    y = CONTENT_TOP + i * Inches(0.95)
    add_card(slide, Inches(7.0), y, Inches(5.65), Inches(0.82), CARD_BG2)
    eb = slide.shapes.add_textbox(Inches(7.15), y + Inches(0.1), Inches(0.9), Inches(0.35))
    ef = eb.text_frame
    ep = ef.paragraphs[0]
    ep.text = ext
    ep.font.size = Pt(12)
    ep.font.bold = True
    ep.font.color.rgb = CYAN
    ep.font.name = FONT_EN
    lb = slide.shapes.add_textbox(Inches(8.15), y + Inches(0.1), Inches(1.8), Inches(0.35))
    lf = lb.text_frame
    lp = lf.paragraphs[0]
    lp.text = lib
    lp.font.size = Pt(12)
    lp.font.bold = True
    lp.font.color.rgb = WHITE
    lp.font.name = FONT
    db = slide.shapes.add_textbox(Inches(7.15), y + Inches(0.45), Inches(5.3), Inches(0.35))
    df = db.text_frame
    dp = df.paragraphs[0]
    dp.text = desc
    dp.font.size = Pt(11)
    dp.font.color.rgb = TEXT_SUB
    dp.font.name = FONT

# 代码块
code = """async function handleKbUpload(file) {
  const formData = new FormData();
  formData.append('file', file);
  const res = await fetch('/api/kb/upload', {
    method: 'POST',
    headers: { 'Authorization': 'Bearer ' + token },
    body: formData
  });
}"""
add_code_block(slide, code, MARGIN_L, Inches(5.0), Inches(12.0), Inches(1.6))

# ============================================================================
# 第19页：Render 部署实战
# ============================================================================
slide = add_page(19, "Render 部署实战", "免费版下的工程化权衡", ORANGE)

# 架构图
add_card(slide, MARGIN_L, CONTENT_TOP, Inches(12.0), Inches(3.2), CARD_BG2)

# 用户图标
user_box = slide.shapes.add_textbox(MARGIN_L + Inches(0.5), CONTENT_TOP + Inches(0.8), Inches(1.0), Inches(0.5))
uf = user_box.text_frame
up = uf.paragraphs[0]
up.text = "用户浏览器"
up.font.size = Pt(12)
up.font.bold = True
up.font.color.rgb = WHITE
up.font.name = FONT
up.alignment = PP_ALIGN.CENTER

# Render Web Service
ws = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L + Inches(2.3), CONTENT_TOP + Inches(0.4), Inches(3.2), Inches(1.3))
ws.fill.solid()
ws.fill.fore_color.rgb = RGBColor(0x1a, 0x3a, 0x6a)
ws.line.color.rgb = CYAN
ws.line.width = Pt(1.5)
ws_tf = ws.text_frame
ws_p = ws_tf.paragraphs[0]
ws_p.text = "Render Web Service"
ws_p.font.size = Pt(15)
ws_p.font.bold = True
ws_p.font.color.rgb = WHITE
ws_p.font.name = FONT
ws_p.alignment = PP_ALIGN.CENTER
ws_p2 = ws_tf.add_paragraph()
ws_p2.text = "Node.js + Express"
ws_p2.font.size = Pt(12)
ws_p2.font.color.rgb = TEXT_SUB
ws_p2.font.name = FONT
ws_p2.alignment = PP_ALIGN.CENTER

# PostgreSQL
pg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L + Inches(6.5), CONTENT_TOP + Inches(0.4), Inches(2.4), Inches(1.3))
pg.fill.solid()
pg.fill.fore_color.rgb = RGBColor(0x1a, 0x4a, 0x4a)
pg.line.color.rgb = GREEN
pg.line.width = Pt(1.5)
pg_tf = pg.text_frame
pg_p = pg_tf.paragraphs[0]
pg_p.text = "PostgreSQL"
pg_p.font.size = Pt(15)
pg_p.font.bold = True
pg_p.font.color.rgb = WHITE
pg_p.font.name = FONT
pg_p.alignment = PP_ALIGN.CENTER
pg_p2 = pg_tf.add_paragraph()
pg_p2.text = "Render.com"
pg_p2.font.size = Pt(12)
pg_p2.font.color.rgb = TEXT_SUB
pg_p2.font.name = FONT
pg_p2.alignment = PP_ALIGN.CENTER

# 大模型 API
ai = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, MARGIN_L + Inches(9.6), CONTENT_TOP + Inches(0.4), Inches(2.0), Inches(1.3))
ai.fill.solid()
ai.fill.fore_color.rgb = RGBColor(0x3a, 0x1a, 0x6a)
ai.line.color.rgb = PURPLE
ai.line.width = Pt(1.5)
ai_tf = ai.text_frame
ai_p = ai_tf.paragraphs[0]
ai_p.text = "Moonshot /"
ai_p.font.size = Pt(13)
ai_p.font.bold = True
ai_p.font.color.rgb = WHITE
ai_p.font.name = FONT
ai_p.alignment = PP_ALIGN.CENTER
ai_p2 = ai_tf.add_paragraph()
ai_p2.text = "DeepSeek API"
ai_p2.font.size = Pt(13)
ai_p2.font.bold = True
ai_p2.font.color.rgb = WHITE
ai_p2.font.name = FONT
ai_p2.alignment = PP_ALIGN.CENTER

# 连接线
line1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L + Inches(1.5), CONTENT_TOP + Inches(1.0), Inches(0.8), Inches(0.03))
line1.fill.solid(); line1.fill.fore_color.rgb = CYAN; line1.line.fill.background()
line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L + Inches(5.5), CONTENT_TOP + Inches(1.0), Inches(1.0), Inches(0.03))
line2.fill.solid(); line2.fill.fore_color.rgb = CYAN; line2.line.fill.background()
line3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L + Inches(8.9), CONTENT_TOP + Inches(1.0), Inches(0.7), Inches(0.03))
line3.fill.solid(); line3.fill.fore_color.rgb = PURPLE; line3.line.fill.background()

add_bullets(slide, [
    "注意：免费 Web Service 会在 15 分钟无流量后休眠，首次请求有冷启动延迟",
    "PostgreSQL 免费实例同样会休眠，但数据不会丢失",
    "生产环境建议升级到付费实例或增加健康检查探针"
], MARGIN_L, CONTENT_TOP + Inches(3.45), Inches(12.0), Inches(1.6), BODY_SIZE, TEXT_MAIN)

add_bullets(slide, [
    "render.yaml 定义了环境变量模板，不包含敏感值，便于团队协作"
], MARGIN_L, Inches(6.6), Inches(12.0), Inches(0.5), SMALL_SIZE, TEXT_SUB)

# ============================================================================
# 第20页：踩坑实录（一）
# ============================================================================
slide = add_page(20, "踩坑实录（一）", "持久化与认证", RED)

pitfalls = [
    ("坑 1：Render 本地文件丢失", "免费 Web Service 无 Persistent Disk，SQLite 在实例休眠后清空，导致用户全部消失。", "迁移至 PostgreSQL，所有状态写入外部数据库。", GREEN),
    ("坑 2：JWT Secret 不稳定", "早期用随机字符串，每次部署后所有用户 Token 失效，被迫重新登录。", "将 JWT_SECRET 固定为环境变量，部署时保持不变。", GREEN),
    ("坑 3：密码明文 vs 哈希", "初期为快速验证未做 bcrypt，存在严重安全隐患。", "注册/登录统一使用 bcryptjs 加盐哈希。", GREEN),
]
for i, (title, problem, solution, color) in enumerate(pitfalls):
    y = CONTENT_TOP + i * Inches(1.65)
    add_card(slide, MARGIN_L, y, Inches(12.0), Inches(1.45), CARD_BG2)
    # 标题
    tb = slide.shapes.add_textbox(MARGIN_L + Inches(0.2), y + Inches(0.12), Inches(11.6), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RED
    p.font.name = FONT
    # 问题
    pb = slide.shapes.add_textbox(MARGIN_L + Inches(0.2), y + Inches(0.45), Inches(11.6), Inches(0.45))
    pf = pb.text_frame
    pf.word_wrap = True
    pp = pf.paragraphs[0]
    pp.text = "问题：" + problem
    pp.font.size = SMALL_SIZE
    pp.font.color.rgb = TEXT_SUB
    pp.font.name = FONT
    # 方案
    sb = slide.shapes.add_textbox(MARGIN_L + Inches(0.2), y + Inches(0.95), Inches(11.6), Inches(0.45))
    sf = sb.text_frame
    sf.word_wrap = True
    sp = sf.paragraphs[0]
    sp.text = "方案：" + solution
    sp.font.size = SMALL_SIZE
    sp.font.color.rgb = color
    sp.font.name = FONT

# ============================================================================
# 第21页：踩坑实录（二）
# ============================================================================
slide = add_page(21, "踩坑实录（二）", "模型、检索与交互状态", RED)

pitfalls2 = [
    ("坑 4：模型参数不兼容", "某些模型对 temperature 取值范围敏感，调用报错。", "统一设置 temperature=1，屏蔽模型差异。"),
    ("坑 5：中文检索按空格分词", "RAG 召回率低，AI 回答『我不知道』。", "中文按字拆分，英文按词拆分，召回率显著提升。"),
    ("坑 6：Bond Codin 状态混乱", "普通对话与代码审计模式共用一套状态，切换时历史串台。", "引入 isBondCodinMode + mode='bondcodin' 字段，独立保存与恢复。"),
    ("坑 7：图片 OCR 文本质量差", "PDF 扫描件识别错误导致 RAG 噪声。", "前端限制文件类型 + 后端记录原始文件名便于追溯。"),
]
for i, (title, problem, solution) in enumerate(pitfalls2):
    y = CONTENT_TOP + i * Inches(1.25)
    add_card(slide, MARGIN_L, y, Inches(12.0), Inches(1.1), CARD_BG2)
    # 标题
    tb = slide.shapes.add_textbox(MARGIN_L + Inches(0.2), y + Inches(0.08), Inches(11.6), Inches(0.3))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(15)
    p.font.bold = True
    p.font.color.rgb = ORANGE
    p.font.name = FONT
    # 内容
    pb = slide.shapes.add_textbox(MARGIN_L + Inches(0.2), y + Inches(0.4), Inches(5.6), Inches(0.6))
    pf = pb.text_frame
    pf.word_wrap = True
    pp = pf.paragraphs[0]
    pp.text = "问题：" + problem
    pp.font.size = Pt(11)
    pp.font.color.rgb = TEXT_SUB
    pp.font.name = FONT
    sb = slide.shapes.add_textbox(MARGIN_L + Inches(6.0), y + Inches(0.4), Inches(5.8), Inches(0.6))
    sf = sb.text_frame
    sf.word_wrap = True
    sp = sf.paragraphs[0]
    sp.text = "方案：" + solution
    sp.font.size = Pt(11)
    sp.font.color.rgb = GREEN
    sp.font.name = FONT

# ============================================================================
# 第22页：安全与边界
# ============================================================================
slide = add_page(22, "安全与边界", "AI 应用必须守住的底线", RED)

add_two_col_text(
    slide,
    "已落地的安全措施",
    [
        "JWT + bcrypt：接口全量鉴权，密码加盐哈希",
        "用户数据隔离：所有查询都带 user_id 过滤",
        "环境变量隔离：密钥不上传仓库",
        "文件类型白名单：防止恶意文件上传"
    ],
    "仍需关注的边界",
    [
        "模型幻觉：RAG 只能降低，无法根除",
        "Prompt 注入：用户可能通过输入绕过角色约束",
        "敏感数据：上传文档可能含内部信息，需审计",
        "成本控制：API 调用按 token 计费，需限流"
    ],
    CONTENT_TOP,
    GREEN,
    ORANGE
)

add_bullets(slide, [
    "原则：AI 应用『先可用，再安全』，但安全不能无限期后置；特别是金融领域，合规与数据保护必须尽早纳入设计。"
], MARGIN_L, Inches(5.6), Inches(12.0), Inches(0.8), BODY_SIZE, TEXT_MAIN)

# ============================================================================
# 第23页：成本与性能
# ============================================================================
slide = add_page(23, "成本与性能", "免费版下的现实约束", ORANGE)

add_bullets(slide, [
    "Render 免费版：Web Service + PostgreSQL 均免费，但都有休眠机制",
    "大模型 API：按 token 计费，长对话和大知识库会显著增加成本",
    "冷启动：实例休眠后首次请求延迟可达 10-30 秒",
    "前端 SSE：流式输出提升体验，但不减少后端总 token 成本"
], MARGIN_L, CONTENT_TOP, Inches(12.0), Inches(1.7))

# 优化建议卡片
opts = [
    ("会话缓存", "历史对话摘要，减少重复上下文 token"),
    ("知识库切片", "只召回最相关的 Top-K 片段，降低 prompt 长度"),
    ("心跳保活", "定时 ping  health 接口，延缓实例休眠"),
    ("降级策略", "模型不可用时给出友好提示，而非直接报错"),
]
for i, (otitle, odesc) in enumerate(opts):
    x = MARGIN_L + (i % 2) * Inches(6.25)
    y = CONTENT_TOP + Inches(2.1) + (i // 2) * Inches(1.45)
    add_card(slide, x, y, Inches(6.0), Inches(1.25), CARD_BG2)
    tb = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.15), Inches(5.6), Inches(0.35))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    p.text = otitle
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = CYAN
    p.font.name = FONT
    db = slide.shapes.add_textbox(x + Inches(0.2), y + Inches(0.55), Inches(5.6), Inches(0.6))
    df = db.text_frame
    df.word_wrap = True
    dp = df.paragraphs[0]
    dp.text = odesc
    dp.font.size = SMALL_SIZE
    dp.font.color.rgb = TEXT_SUB
    dp.font.name = FONT


# ============================================================================
# 第24页：给其他团队的建议
# ============================================================================
slide = add_page(24, "给其他团队的建议", "如何让 Vibe Coding 真正产生价值", ORANGE)

add_bullets(slide, [
    "从小处着手：先用 AI 生成工具脚本、单元测试、文档，建立信任",
    "保留『人的审核』：AI 生成代码后必须经过 review，尤其是安全、事务、并发相关代码",
    "做好上下文管理：给 AI 的 prompt 要包含足够的业务背景、代码路径、依赖关系",
    "拆分任务：把大需求拆成多个小任务，每个任务独立验证，降低一次性返工风险",
    "记录踩坑：把模型行为、部署问题、业务规则沉淀为团队知识库"
], MARGIN_L, CONTENT_TOP, Inches(12.0), Inches(2.5))

# 建议矩阵
matrix = [
    ("适合 Vibe Coding", "不适合 Vibe Coding"),
    ("原型/MVP、工具脚本、标准 CRUD、文档生成、测试用例", "核心交易逻辑、强合规模块、复杂并发控制、未稳定的需求"),
]
# 表头
add_card(slide, MARGIN_L, CONTENT_TOP + Inches(2.9), Inches(6.0), Inches(0.6), RGBColor(0x1a, 0x4a, 0x4a))
add_card(slide, Inches(6.9), CONTENT_TOP + Inches(2.9), Inches(6.0), Inches(0.6), RGBColor(0x4a, 0x1a, 0x1a))
for i, header in enumerate(matrix[0]):
    x = MARGIN_L + i * Inches(6.25)
    hb = slide.shapes.add_textbox(x + Inches(0.2), CONTENT_TOP + Inches(3.0), Inches(5.6), Inches(0.4))
    hf = hb.text_frame
    hp = hf.paragraphs[0]
    hp.text = header
    hp.font.size = Pt(14)
    hp.font.bold = True
    hp.font.color.rgb = WHITE
    hp.font.name = FONT
    hp.alignment = PP_ALIGN.CENTER
# 内容行
add_card(slide, MARGIN_L, CONTENT_TOP + Inches(3.55), Inches(6.0), Inches(1.1), CARD_BG2)
add_card(slide, Inches(6.9), CONTENT_TOP + Inches(3.55), Inches(6.0), Inches(1.1), CARD_BG2)
for i, content in enumerate(matrix[1]):
    x = MARGIN_L + i * Inches(6.25)
    cb = slide.shapes.add_textbox(x + Inches(0.25), CONTENT_TOP + Inches(3.7), Inches(5.5), Inches(0.9))
    cf = cb.text_frame
    cf.word_wrap = True
    cp = cf.paragraphs[0]
    cp.text = content
    cp.font.size = SMALL_SIZE
    cp.font.color.rgb = TEXT_MAIN if i == 0 else TEXT_SUB
    cp.font.name = FONT

add_bullets(slide, [
    "核心认知：Vibe Coding 是加速器，不是替代品；最终决定质量的还是开发者的判断力和架构能力。"
], MARGIN_L, Inches(6.5), Inches(12.0), Inches(0.6), BODY_SIZE, TEXT_MAIN)

# ============================================================================
# 第25页：未来演进
# ============================================================================
slide = add_page(25, "未来演进", "从单一 Agent 到智能体协作", CYAN)

add_bullets(slide, [
    "MCP（Model Context Protocol）：标准化模型与外部工具/数据的交互协议",
    "多 Agent 协作：不同领域的专家 Agent 分工完成复杂任务，例如研究 Agent + 审计 Agent + 报告 Agent",
    "本地模型 + 云端混合：敏感代码/数据本地化推理，通用问答调用云端大模型",
    "向量数据库：从关键词检索升级到语义检索，进一步提升 RAG 效果",
    "AI 辅助编码工作流：从 Copilot 建议到自动化测试、文档、部署的全链路"
], MARGIN_L, CONTENT_TOP, Inches(12.0), Inches(2.5))

# 演进时间线
evolve = [
    ("当前", "单 Agent · 关键词 RAG · Render 部署"),
    ("近期", "多模型路由 · 向量检索 · 会话摘要"),
    ("中期", "多 Agent 协作 · MCP 工具链 · 本地模型"),
    ("远期", "企业知识中枢 · 自动化工作流 · 智能决策辅助"),
]
line_y = Inches(5.5)
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGIN_L, line_y + Inches(0.25), Inches(12.0), Inches(0.04))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(0x24, 0x35, 0x5c)
line.line.fill.background()
for i, (stage, desc) in enumerate(evolve):
    x = MARGIN_L + i * Inches(3.05)
    node = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.5), line_y, Inches(0.5), Inches(0.5))
    node.fill.solid()
    node.fill.fore_color.rgb = CYAN if i == 0 else PURPLE if i == 1 else ORANGE if i == 2 else GREEN
    node.line.fill.background()
    # stage
    sb = slide.shapes.add_textbox(x, line_y + Inches(0.65), Inches(1.5), Inches(0.3))
    sf = sb.text_frame
    sp = sf.paragraphs[0]
    sp.text = stage
    sp.font.size = Pt(13)
    sp.font.bold = True
    sp.font.color.rgb = WHITE
    sp.font.name = FONT
    sp.alignment = PP_ALIGN.CENTER
    # desc
    db = slide.shapes.add_textbox(x - Inches(0.3), line_y + Inches(1.0), Inches(2.1), Inches(0.8))
    df = db.text_frame
    df.word_wrap = True
    dp = df.paragraphs[0]
    dp.text = desc
    dp.font.size = Pt(10)
    dp.font.color.rgb = TEXT_SUB
    dp.font.name = FONT
    dp.alignment = PP_ALIGN.CENTER

# ============================================================================
# 第26页：总结
# ============================================================================
slide = add_page(26, "总结", "把 AI 当作放大器", CYAN)

add_bullets(slide, [
    "Vibe Coding 不是不写代码，而是用自然语言驱动 AI 生成代码，开发者聚焦架构与质量",
    "ccdc-ai-service 验证了一个完整 AI 应用可以在极少人力下快速落地：Express + 零框架前端 + PostgreSQL + 多模型",
    "Agent 设计的关键：清晰的角色、可控的工具、结构化的输出、人的持续监督",
    "企业级落地仍需关注：数据持久化、安全认证、模型边界、幻觉与成本控制",
    "最重要的不是追逐热点，而是把 AI 嵌入到真实工作流中，解决具体问题"
], MARGIN_L, CONTENT_TOP, Inches(12.0), Inches(3.2))

# 金句卡片
add_card(slide, MARGIN_L, Inches(5.6), Inches(12.0), Inches(1.1), RGBColor(0x1a, 0x3a, 0x5a))
quote_box = slide.shapes.add_textbox(MARGIN_L + Inches(0.3), Inches(5.85), Inches(11.4), Inches(0.7))
qf = quote_box.text_frame
qp = qf.paragraphs[0]
qp.text = "\"AI 不会取代程序员，但会用 AI 的程序员会取代不会用的。\""
qp.font.size = Pt(20)
qp.font.bold = True
qp.font.color.rgb = CYAN
qp.font.name = FONT
qp.alignment = PP_ALIGN.CENTER

# ============================================================================
# 第27页：Q&A
# ============================================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_background(slide)
add_top_bar(slide, CYAN)
add_glow_circle(slide, Inches(4.5), Inches(1.8), Inches(4.5), Inches(4.5), CYAN, 0.9)
add_glow_circle(slide, Inches(7.5), Inches(2.2), Inches(3.0), Inches(3.0), PURPLE, 0.92)

title_box = slide.shapes.add_textbox(MARGIN_L, Inches(2.8), Inches(12.0), Inches(1.5))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Q & A"
p.font.size = Pt(66)
p.font.bold = True
p.font.color.rgb = WHITE
p.font.name = FONT
p.alignment = PP_ALIGN.CENTER
p2 = tf.add_paragraph()
p2.text = "欢迎交流与探讨"
p2.font.size = Pt(22)
p2.font.color.rgb = CYAN
p2.font.name = FONT
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(14)

add_footer(slide, 27)

# ============================================================================
# 保存
# ============================================================================
output_path = "CCDC_AI_VibeCoding_Training_v2.pptx"
prs.save(output_path)
print(f"PPT 已生成：{output_path}")
print(f"总页数：{len(prs.slides)}")
