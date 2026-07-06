# -*- coding: utf-8 -*-
"""
仅替换 CCDC_AI_VibeCoding_Training_v2.pptx 的背景，不修改内容。
新背景：渐变底色 + 科技网格 + 半透明光晕。
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
import copy

INPUT_PATH = "CCDC_AI_VibeCoding_Training_v2.pptx"
OUTPUT_PATH = "CCDC_AI_VibeCoding_Training_v3.pptx"

prs = Presentation(INPUT_PATH)
W = prs.slide_width
H = prs.slide_height

# 新背景配色
BG_TOP = RGBColor(0x0a, 0x16, 0x3d)      # 左上：深蓝（提亮，避免纯黑）
BG_BOTTOM = RGBColor(0x2d, 0x10, 0x52)   # 右下：深紫（提亮，增强渐变感）
GRID_COLOR = RGBColor(0x2d, 0x4a, 0x80)  # 网格线（稍亮，增加层次）
GLOW_CYAN = RGBColor(0x00, 0xe5, 0xff)   # 青色光晕
GLOW_PURPLE = RGBColor(0x8b, 0x5c, 0xff) # 紫色光晕


def is_fullscreen_rect(shape):
    """判断是否全页背景矩形（允许微小偏移）"""
    try:
        return (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and
                abs(shape.left) < Inches(0.3) and abs(shape.top) < Inches(0.3) and
                abs(shape.width - W) < Inches(0.3) and abs(shape.height - H) < Inches(0.3))
    except Exception:
        return False


def is_top_bar(shape):
    """判断是否顶部装饰条"""
    try:
        return (shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and
                abs(shape.left) < Inches(0.3) and abs(shape.top) < Inches(0.3) and
                abs(shape.width - W) < Inches(0.3) and shape.height < Inches(0.15))
    except Exception:
        return False


def is_large_glow_circle(shape):
    """判断是否大光晕圆（尺寸大、纯色、颜色为常见光晕色）"""
    try:
        if shape.shape_type != MSO_SHAPE_TYPE.AUTO_SHAPE:
            return False
        aspect = shape.width / shape.height
        if not (0.85 < aspect < 1.15):
            return False
        if shape.width < Inches(2.0):
            return False
        fill = shape.fill
        if fill.type != 1:  # solid
            return False
        if fill.fore_color.type is None:
            return False
        rgb = str(fill.fore_color.rgb).upper()
        # 常见光晕色：青色、紫色、蓝色、白色
        return rgb in ['00E5FF', '8B5CFF', '4A90D9', 'FFFFFF']
    except Exception:
        return False


def move_to_bottom(shape):
    """将形状移到最底层（spTree 中索引 2 的位置）"""
    sp = shape._element
    spTree = sp.getparent()
    spTree.remove(sp)
    spTree.insert(2, sp)


def remove_old_background(slide):
    """删除旧的背景元素：全页矩形、顶部条、大光晕圆"""
    to_remove = []
    for shape in slide.shapes:
        if is_fullscreen_rect(shape) or is_top_bar(shape) or is_large_glow_circle(shape):
            to_remove.append(shape)
    print(f"  发现 {len(to_remove)} 个旧背景元素")
    for shape in to_remove:
        try:
            sp = shape._element
            sp.getparent().remove(sp)
            print(f"    已删除: {shape.shape_type} {shape.width/914400:.2f}x{shape.height/914400:.2f}")
        except Exception as e:
            print(f"  删除背景元素失败: {e}")


def add_gradient_background(slide, angle=315):
    """添加渐变全页背景"""
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    fill = bg.fill
    fill.gradient()
    fill.gradient_angle = angle
    # 清空默认停止点并重新添加
    gs = fill.gradient_stops
    # python-pptx 默认会创建 2 个停止点
    stops = list(gs)
    if len(stops) >= 2:
        stops[0].position = 0.0
        stops[0].color.rgb = BG_TOP
        stops[1].position = 1.0
        stops[1].color.rgb = BG_BOTTOM
    bg.line.fill.background()
    move_to_bottom(bg)
    return bg


def add_grid(slide, spacing=Inches(1.2), color=GRID_COLOR, transparency=0.75):
    """添加细密科技网格线"""
    # 横线
    y = spacing
    while y < H:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, y, W, Pt(0.5))
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.fill.transparency = transparency
        line.line.fill.background()
        move_to_bottom(line)
        y += spacing
    # 竖线
    x = spacing
    while x < W:
        line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, 0, Pt(0.5), H)
        line.fill.solid()
        line.fill.fore_color.rgb = color
        line.fill.transparency = transparency
        line.line.fill.background()
        move_to_bottom(line)
        x += spacing


def add_glow_circle(slide, x, y, size, color, transparency=0.92):
    """添加半透明光晕圆"""
    c = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, size, size)
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.fill.transparency = transparency
    c.line.fill.background()
    move_to_bottom(c)
    return c


def add_top_bar(slide, color=RGBColor(0x00, 0xe5, 0xff)):
    """添加顶部细装饰条"""
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    move_to_bottom(bar)


def add_vignette(slide):
    """添加四角暗角，增强视觉聚焦"""
    # 上方暗角
    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Inches(1.8))
    fill = top.fill
    fill.gradient()
    fill.gradient_angle = 270
    stops = list(fill.gradient_stops)
    if len(stops) >= 2:
        stops[0].position = 0.0
        stops[0].color.rgb = RGBColor(0x06, 0x0c, 0x28)
        stops[1].position = 1.0
        stops[1].color.rgb = RGBColor(0x06, 0x0c, 0x28)
    top.fill.transparency = 0.4
    top.line.fill.background()
    move_to_bottom(top)

    # 左下暗角
    left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, H - Inches(2.0), Inches(2.0), Inches(2.0))
    fill = left.fill
    fill.gradient()
    fill.gradient_angle = 0
    stops = list(fill.gradient_stops)
    if len(stops) >= 2:
        stops[0].position = 0.0
        stops[0].color.rgb = RGBColor(0x06, 0x0c, 0x28)
        stops[1].position = 1.0
        stops[1].color.rgb = RGBColor(0x06, 0x0c, 0x28)
    left.fill.transparency = 0.4
    left.line.fill.background()
    move_to_bottom(left)


# 处理每一页
for idx, slide in enumerate(prs.slides, 1):
    print(f"处理第 {idx} 页...")
    remove_old_background(slide)
    # 注意：越先 move_to_bottom 的元素在背景层越靠上；渐变背景必须最后添加，才能处于最底层
    # 1. 顶部装饰条
    add_top_bar(slide, RGBColor(0x4a, 0x90, 0xd9))
    # 2. 暗角
    add_vignette(slide)
    # 3. 光晕（根据页码变化，避免单调）
    add_glow_circle(slide, W - Inches(3.8), Inches(-1.5), Inches(5.0), GLOW_CYAN, 0.93)
    if idx % 2 == 0:
        add_glow_circle(slide, Inches(-1.5), H - Inches(3.0), Inches(4.0), GLOW_PURPLE, 0.94)
    if idx % 3 == 0:
        add_glow_circle(slide, W / 2 - Inches(1.5), H - Inches(1.5), Inches(3.0), GLOW_PURPLE, 0.95)
    # 4. 科技网格
    add_grid(slide)
    # 5. 渐变背景（最后添加，确保在最底层）
    add_gradient_background(slide)

prs.save(OUTPUT_PATH)
print(f"\n新 PPT 已保存：{OUTPUT_PATH}")
print(f"总页数：{len(prs.slides)}")
